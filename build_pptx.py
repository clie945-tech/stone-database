"""
ReStone 循環石材服務平台 — 中英雙語簡介簡報
產出：ReStone_簡介簡報.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ─── 品牌色系 ───
NAVY_DEEP  = RGBColor(0x0A, 0x0E, 0x1A)
NAVY_MID   = RGBColor(0x1E, 0x1B, 0x4B)
BLUE       = RGBColor(0x25, 0x63, 0xEB)
BLUE_LIGHT = RGBColor(0x60, 0xA5, 0xFA)
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_LIGHT = RGBColor(0xA7, 0x8B, 0xFA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF8, 0xFA, 0xFC)
TEXT_DARK  = RGBColor(0x0F, 0x17, 0x2A)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
TEXT_LIGHT = RGBColor(0x94, 0xA3, 0xB8)
BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
GREEN      = RGBColor(0x10, 0xB9, 0x81)
GOLD       = RGBColor(0xF5, 0x9E, 0x0B)
PINK       = RGBColor(0xEC, 0x48, 0x99)

# 中文字型優先（PPT 在 Windows / Mac 都有）
FONT_HEAD = 'Microsoft JhengHei'   # 標題
FONT_BODY = 'Microsoft JhengHei'   # 內文
FONT_EN   = 'Calibri'              # 英文副標

# ─── 建立簡報（16:9）───
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank = prs.slide_layouts[6]  # 完全空白版型


# ─── 工具函式 ───
def add_rect(slide, x, y, w, h, fill, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w: s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_rounded(slide, x, y, w, h, fill, line=None, line_w=None, radius=0.06):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w: s.line.width = line_w
    s.shadow.inherit = False
    # 調整圓角
    try:
        s.adjustments[0] = radius
    except Exception:
        pass
    return s


def add_hex(slide, x, y, w, h, fill, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w: s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_oval(slide, x, y, w, h, fill, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w: s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TEXT_DARK,
             align='left', valign='top', font=None, italic=False, spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.margin_left = Emu(0)
    tb.text_frame.margin_right = Emu(0)
    tb.text_frame.margin_top = Emu(0)
    tb.text_frame.margin_bottom = Emu(0)
    tb.text_frame.word_wrap = True
    valign_map = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}
    tb.text_frame.vertical_anchor = valign_map.get(valign, MSO_ANCHOR.TOP)
    align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
    p = tb.text_frame.paragraphs[0]
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if font: r.font.name = font
    if spacing is not None:
        p.line_spacing = spacing
    return tb


def add_multiline(slide, x, y, w, h, lines, *, align='left', valign='top'):
    """每行可獨立指定樣式 [(text, opts), ...] opts 含 size/bold/color/font"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    valign_map = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = valign_map.get(valign, MSO_ANCHOR.TOP)
    align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
    for i, (text, opts) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)
        if 'space_before' in opts:
            p.space_before = Pt(opts['space_before'])
        if 'space_after' in opts:
            p.space_after = Pt(opts['space_after'])
        r = p.add_run()
        r.text = text
        r.font.size = Pt(opts.get('size', 14))
        r.font.bold = opts.get('bold', False)
        r.font.italic = opts.get('italic', False)
        r.font.color.rgb = opts.get('color', TEXT_DARK)
        if 'font' in opts:
            r.font.name = opts['font']
    return tb


def draw_diamond_logo(slide, cx_in, cy_in, size_in, accent=BLUE, accent2=PURPLE):
    """繪製簡化版切面鑽石 LOGO（六角形 + 中央菱形）"""
    s = size_in
    x = cx_in - s / 2
    y = cy_in - s / 2
    # 外框六角形
    hex_shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(x), Inches(y), Inches(s), Inches(s))
    hex_shape.fill.solid(); hex_shape.fill.fore_color.rgb = accent
    hex_shape.line.fill.background()
    hex_shape.shadow.inherit = False
    hex_shape.rotation = 90  # 旋轉成立式六角形
    # 中央菱形
    inner_s = s * 0.4
    ix = cx_in - inner_s / 2
    iy = cy_in - inner_s / 2
    diamond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(ix), Inches(iy), Inches(inner_s), Inches(inner_s))
    diamond.fill.solid(); diamond.fill.fore_color.rgb = accent2
    diamond.line.fill.background()
    diamond.shadow.inherit = False


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 1: 封面 Title Hero
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, SW, SH, NAVY_DEEP)
# 背景裝飾：左下角大鑽石（淡）
add_hex(slide, Inches(-2), Inches(4.5), Inches(8), Inches(8), NAVY_MID)
# Logo
draw_diamond_logo(slide, 2.5, 2.6, 1.4, BLUE, PURPLE_LIGHT)
# 標題
add_text(slide, Inches(4.2), Inches(2.0), Inches(8), Inches(1.2), 'ReStone',
         size=72, bold=True, color=WHITE, font=FONT_EN)
# 副標
add_text(slide, Inches(4.2), Inches(3.0), Inches(8), Inches(0.6), '循環石材服務平台',
         size=28, bold=True, color=BLUE_LIGHT, font=FONT_HEAD)
add_text(slide, Inches(4.2), Inches(3.55), Inches(8), Inches(0.4), 'Circular Stone Service Platform',
         size=14, color=TEXT_LIGHT, font=FONT_EN)
# 主訴求
add_text(slide, Inches(4.2), Inches(4.3), Inches(8), Inches(0.6),
         'AI 驅動的循環石材設計平台',
         size=16, color=WHITE, font=FONT_HEAD)
add_text(slide, Inches(4.2), Inches(4.75), Inches(8), Inches(0.4),
         'AI-Powered Design Platform for Sustainable Stone',
         size=11, color=TEXT_LIGHT, italic=True, font=FONT_EN)
# 底部資訊
add_text(slide, Inches(0.6), Inches(7), Inches(8), Inches(0.3),
         '功能簡介簡報  ·  Product Introduction',
         size=10, color=TEXT_LIGHT, font=FONT_EN)
add_text(slide, Inches(10), Inches(7), Inches(3), Inches(0.3),
         '2026 · v1.0',
         size=10, color=TEXT_LIGHT, align='right', font=FONT_EN)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 2: 議程 Agenda
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, SW, SH, WHITE)
# 標題
add_text(slide, Inches(0.7), Inches(0.6), Inches(12), Inches(0.6), '簡報大綱',
         size=36, bold=True, color=TEXT_DARK, font=FONT_HEAD)
add_text(slide, Inches(0.7), Inches(1.2), Inches(12), Inches(0.3), 'AGENDA',
         size=11, color=PURPLE, bold=True, font=FONT_EN)

agenda_items = [
    ('01', '平台願景', 'Platform Vision'),
    ('02', '核心功能', 'Core Features'),
    ('03', 'AI 旗艦：設計合成', 'Flagship: AI Design'),
    ('04', '六大 AI 引擎', '6 AI Engines'),
    ('05', '使用者角色', 'User Roles'),
    ('06', '技術架構', 'Tech Stack'),
    ('07', '資安合規', 'Security & Compliance'),
    ('08', '成本結構', 'Cost Structure'),
]

cols = 4; rows = 2
gap_x = 0.25; gap_y = 0.35
card_w = (13.333 - 1.4 - gap_x * (cols - 1)) / cols
card_h = 1.95
start_x = 0.7
start_y = 2.0

for i, (num, zh, en) in enumerate(agenda_items):
    col = i % cols
    row = i // cols
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    add_rounded(slide, Inches(x), Inches(y), Inches(card_w), Inches(card_h),
                LIGHT_BG, line=BORDER, line_w=Pt(0.75), radius=0.08)
    # 編號圓點
    add_oval(slide, Inches(x + 0.3), Inches(y + 0.3), Inches(0.55), Inches(0.55),
             BLUE)
    add_text(slide, Inches(x + 0.3), Inches(y + 0.36), Inches(0.55), Inches(0.45),
             num, size=14, bold=True, color=WHITE, align='center', font=FONT_EN)
    # 中文名
    add_text(slide, Inches(x + 0.3), Inches(y + 1.0), Inches(card_w - 0.6), Inches(0.4),
             zh, size=16, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    # 英文名
    add_text(slide, Inches(x + 0.3), Inches(y + 1.42), Inches(card_w - 0.6), Inches(0.35),
             en, size=10, color=TEXT_MUTED, font=FONT_EN)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 3: 平台願景 Platform Vision
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, SW, SH, WHITE)

add_text(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4), 'SECTION 01',
         size=11, bold=True, color=PURPLE, font=FONT_EN)
add_text(slide, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         '從找料目錄，升級為 AI 設計工具',
         size=32, bold=True, color=TEXT_DARK, font=FONT_HEAD)
add_text(slide, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
         'From Catalog to AI-Powered Design Tool',
         size=14, color=TEXT_MUTED, italic=True, font=FONT_EN)

# 左：痛點
add_rounded(slide, Inches(0.7), Inches(2.5), Inches(5.8), Inches(4.4),
            RGBColor(0xFE, 0xF2, 0xF2), line=RGBColor(0xFE, 0xCA, 0xCA), line_w=Pt(0.75), radius=0.05)
add_text(slide, Inches(1.0), Inches(2.8), Inches(5.2), Inches(0.4), '⚠ 業界痛點 / Pain Points',
         size=14, bold=True, color=RGBColor(0xB9, 0x1C, 0x1C), font=FONT_HEAD)
problems = [
    ('提案改了七版', 'Seven revisions, still uncertain'),
    ('客戶看不懂規格表', 'Spec sheets don\'t sell to clients'),
    ('每張合成圖 30 分鐘 PS', '30 mins per Photoshop edit'),
    ('剩料、邊角料賣不掉', 'Circular stones sit unsold'),
]
for i, (zh, en) in enumerate(problems):
    y = 3.4 + i * 0.85
    add_text(slide, Inches(1.0), Inches(y), Inches(0.3), Inches(0.4), '✕',
             size=14, bold=True, color=RGBColor(0xB9, 0x1C, 0x1C), font=FONT_EN)
    add_text(slide, Inches(1.35), Inches(y - 0.05), Inches(5), Inches(0.4),
             zh, size=14, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    add_text(slide, Inches(1.35), Inches(y + 0.3), Inches(5), Inches(0.3),
             en, size=10, color=TEXT_MUTED, italic=True, font=FONT_EN)

# 右：解決方案
add_rounded(slide, Inches(6.8), Inches(2.5), Inches(5.8), Inches(4.4),
            RGBColor(0xEF, 0xF6, 0xFF), line=RGBColor(0xBF, 0xDB, 0xFE), line_w=Pt(0.75), radius=0.05)
add_text(slide, Inches(7.1), Inches(2.8), Inches(5.2), Inches(0.4), '✓ ReStone 解法 / Our Solution',
         size=14, bold=True, color=BLUE, font=FONT_HEAD)
solutions = [
    ('30 秒 AI 合成', '30-second AI synthesis'),
    ('即時視覺化說服客戶', 'Instant visual proposals'),
    ('一鍵套用，無需軟體', 'One-click, no software needed'),
    ('剩料精準曝光給對的設計', 'Right stones, right designs'),
]
for i, (zh, en) in enumerate(solutions):
    y = 3.4 + i * 0.85
    add_text(slide, Inches(7.1), Inches(y), Inches(0.3), Inches(0.4), '✓',
             size=14, bold=True, color=BLUE, font=FONT_EN)
    add_text(slide, Inches(7.45), Inches(y - 0.05), Inches(5), Inches(0.4),
             zh, size=14, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    add_text(slide, Inches(7.45), Inches(y + 0.3), Inches(5), Inches(0.3),
             en, size=10, color=TEXT_MUTED, italic=True, font=FONT_EN)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 4: 核心功能 Core Features (六大 AI 引擎 grid)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, SW, SH, WHITE)

add_text(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4), 'SECTION 02',
         size=11, bold=True, color=PURPLE, font=FONT_EN)
add_text(slide, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         '六大 AI 引擎，深度整合',
         size=32, bold=True, color=TEXT_DARK, font=FONT_HEAD)
add_text(slide, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
         'Six AI Engines, Deeply Embedded',
         size=14, color=TEXT_MUTED, italic=True, font=FONT_EN)

features = [
    ('🎨', 'AI 設計合成', 'AI Design Synthesis', 'Flux Fill Pro inpainting，30 秒套上石材紋理'),
    ('👁', '視覺分析', 'Vision Analysis', 'LLaVA 看石材照片，自動寫精準描述'),
    ('🔍', '以圖搜圖', 'Image Search', 'pHash + RGB 雙軸比對，秒找相似石材'),
    ('🌈', '顏色偵測', 'Color Detection', 'HSV 色彩空間自動分類 8 種顏色標籤'),
    ('🧠', 'Prompt 工程', 'Prompt Engineering', '三層 fallback：人工 → 視覺 AI → 模板'),
    ('💬', '語義搜尋', 'Semantic Search', '中英同義詞自動展開，跨語言匹配'),
]

cols = 3; rows = 2
gap = 0.25
card_w = (13.333 - 1.4 - gap * (cols - 1)) / cols
card_h = 2.4
start_x = 0.7
start_y = 2.3

for i, (emoji, zh, en, desc) in enumerate(features):
    col = i % cols
    row = i // cols
    x = start_x + col * (card_w + gap)
    y = start_y + row * (card_h + gap)
    add_rounded(slide, Inches(x), Inches(y), Inches(card_w), Inches(card_h),
                LIGHT_BG, line=BORDER, line_w=Pt(0.75), radius=0.04)
    # Icon 圓底
    add_oval(slide, Inches(x + 0.35), Inches(y + 0.35), Inches(0.7), Inches(0.7), BLUE)
    add_text(slide, Inches(x + 0.35), Inches(y + 0.42), Inches(0.7), Inches(0.55),
             emoji, size=22, align='center', color=WHITE, font=FONT_HEAD)
    # 中文名
    add_text(slide, Inches(x + 0.35), Inches(y + 1.2), Inches(card_w - 0.7), Inches(0.4),
             zh, size=15, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    # 英文
    add_text(slide, Inches(x + 0.35), Inches(y + 1.55), Inches(card_w - 0.7), Inches(0.3),
             en, size=10, color=PURPLE, bold=True, font=FONT_EN)
    # 描述
    add_text(slide, Inches(x + 0.35), Inches(y + 1.85), Inches(card_w - 0.7), Inches(0.5),
             desc, size=10, color=TEXT_MUTED, font=FONT_HEAD)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 5: AI 旗艦 — 設計合成（深色 hero）
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, SW, SH, NAVY_DEEP)
# 裝飾鑽石
add_hex(slide, Inches(10), Inches(-2), Inches(8), Inches(8), NAVY_MID)
add_hex(slide, Inches(-3), Inches(5), Inches(6), Inches(6), NAVY_MID)

# 旗艦徽章
add_rounded(slide, Inches(0.7), Inches(0.6), Inches(2.5), Inches(0.4),
            GOLD, radius=0.5)
add_text(slide, Inches(0.7), Inches(0.62), Inches(2.5), Inches(0.4),
         '⭐ FLAGSHIP FEATURE', size=11, bold=True, color=WHITE, align='center', font=FONT_EN)

add_text(slide, Inches(0.7), Inches(1.25), Inches(12), Inches(0.7),
         'AI 設計合成工具',
         size=40, bold=True, color=WHITE, font=FONT_HEAD)
add_text(slide, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
         'AI Design Synthesis Tool',
         size=20, color=BLUE_LIGHT, italic=True, font=FONT_EN)
add_text(slide, Inches(0.7), Inches(2.65), Inches(12), Inches(0.4),
         '30 秒，把任何石材紋理套上你的設計圖',
         size=15, color=WHITE, font=FONT_HEAD)
add_text(slide, Inches(0.7), Inches(3.05), Inches(12), Inches(0.4),
         'Apply any stone texture to your design in 30 seconds',
         size=11, color=TEXT_LIGHT, italic=True, font=FONT_EN)

# 三步驟流程
steps = [
    ('1', '上傳設計圖', 'Upload Design', '室內、家具、產品皆可'),
    ('2', '筆刷圈選區域', 'Mask Surface', '牆面、桌面、地板'),
    ('3', 'AI 合成', 'AI Generate', '30 秒商業級成圖'),
]
step_w = 3.6
gap_w = 0.3
total_w = step_w * 3 + gap_w * 2
start_x = (13.333 - total_w) / 2
step_y = 4.2

for i, (num, zh, en, desc) in enumerate(steps):
    x = start_x + i * (step_w + gap_w)
    add_rounded(slide, Inches(x), Inches(step_y), Inches(step_w), Inches(2.2),
                NAVY_MID, line=PURPLE_LIGHT, line_w=Pt(0.5), radius=0.06)
    # 大號碼
    add_oval(slide, Inches(x + 0.3), Inches(step_y + 0.3), Inches(0.65), Inches(0.65), PURPLE)
    add_text(slide, Inches(x + 0.3), Inches(step_y + 0.38), Inches(0.65), Inches(0.55),
             num, size=16, bold=True, color=WHITE, align='center', font=FONT_EN)
    # 中文步驟名
    add_text(slide, Inches(x + 0.3), Inches(step_y + 1.05), Inches(step_w - 0.6), Inches(0.4),
             zh, size=16, bold=True, color=WHITE, font=FONT_HEAD)
    # 英文
    add_text(slide, Inches(x + 0.3), Inches(step_y + 1.45), Inches(step_w - 0.6), Inches(0.3),
             en, size=10, color=BLUE_LIGHT, bold=True, font=FONT_EN)
    # 描述
    add_text(slide, Inches(x + 0.3), Inches(step_y + 1.75), Inches(step_w - 0.6), Inches(0.35),
             desc, size=10, color=TEXT_LIGHT, font=FONT_HEAD)

# 底部技術註
add_text(slide, Inches(0.7), Inches(6.85), Inches(12), Inches(0.35),
         'POWERED BY  ·  Replicate Flux Fill Pro  ·  Member-only feature  ·  NT$1.5 per generation',
         size=10, color=TEXT_LIGHT, align='center', font=FONT_EN)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 6: 六大 AI 引擎詳細 - 與 slide 4 不同，這頁強調技術
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, SW, SH, WHITE)

add_text(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4), 'SECTION 04',
         size=11, bold=True, color=PURPLE, font=FONT_EN)
add_text(slide, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         'AI 技術骨幹',
         size=32, bold=True, color=TEXT_DARK, font=FONT_HEAD)
add_text(slide, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
         'AI Technology Stack — Cost-Effective Architecture',
         size=14, color=TEXT_MUTED, italic=True, font=FONT_EN)

# 兩欄：AI 模型 vs 處理流程
# 左欄：使用的 AI 服務
add_text(slide, Inches(0.7), Inches(2.2), Inches(6), Inches(0.35),
         '🤖 AI 服務 / AI Services',
         size=14, bold=True, color=BLUE, font=FONT_HEAD)
ai_services = [
    ('Flux Fill Pro', 'Black Forest Labs', '設計合成 inpainting', '~NT$1.5/張'),
    ('LLaVA v1.6', 'Microsoft Research', '視覺看圖寫描述', '~NT$0.5/張'),
    ('Pillow + ImageHash', 'Open Source', '以圖搜圖、顏色偵測', '免費 / Free'),
]
for i, (name, vendor, use, cost) in enumerate(ai_services):
    y = 2.7 + i * 1.15
    add_rounded(slide, Inches(0.7), Inches(y), Inches(5.8), Inches(1.0),
                LIGHT_BG, line=BORDER, line_w=Pt(0.5), radius=0.04)
    add_text(slide, Inches(0.95), Inches(y + 0.15), Inches(4), Inches(0.35),
             name, size=13, bold=True, color=TEXT_DARK, font=FONT_EN)
    add_text(slide, Inches(0.95), Inches(y + 0.5), Inches(4), Inches(0.3),
             f'{vendor} · {use}', size=9.5, color=TEXT_MUTED, font=FONT_EN)
    add_rounded(slide, Inches(5.0), Inches(y + 0.3), Inches(1.4), Inches(0.4),
                BLUE, radius=0.3)
    add_text(slide, Inches(5.0), Inches(y + 0.32), Inches(1.4), Inches(0.4),
             cost, size=10, bold=True, color=WHITE, align='center', font=FONT_EN)

# 右欄：本機處理（無 API 成本）
add_text(slide, Inches(7.0), Inches(2.2), Inches(6), Inches(0.35),
         '⚙ 本機處理 / Local Processing',
         size=14, bold=True, color=PURPLE, font=FONT_HEAD)
local_features = [
    ('語義搜尋', 'Synonym Match', '純 Python 字典查找，零延遲'),
    ('HEIC → JPG 轉檔', 'Format Normalize', 'iPhone 照片自動相容'),
    ('AI Prompt 三層機制', 'Prompt Fallback', '人工 → 視覺 AI → 模板'),
    ('感知雜湊索引', 'pHash Indexing', '上傳即建索引，無需訓練'),
]
for i, (zh, en, desc) in enumerate(local_features):
    y = 2.7 + i * 0.95
    add_rounded(slide, Inches(7.0), Inches(y), Inches(5.6), Inches(0.85),
                RGBColor(0xF3, 0xE8, 0xFF), line=RGBColor(0xDD, 0xD6, 0xFE), line_w=Pt(0.5), radius=0.04)
    add_text(slide, Inches(7.25), Inches(y + 0.1), Inches(5), Inches(0.35),
             zh, size=12, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    add_text(slide, Inches(7.25), Inches(y + 0.42), Inches(5), Inches(0.25),
             en, size=9.5, color=PURPLE, bold=True, font=FONT_EN)
    add_text(slide, Inches(7.25), Inches(y + 0.62), Inches(5), Inches(0.25),
             desc, size=9, color=TEXT_MUTED, font=FONT_HEAD)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 7: 三種使用者角色 / 3 User Roles
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, SW, SH, WHITE)

add_text(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4), 'SECTION 05',
         size=11, bold=True, color=PURPLE, font=FONT_EN)
add_text(slide, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         '三種角色，分層權限',
         size=32, bold=True, color=TEXT_DARK, font=FONT_HEAD)
add_text(slide, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
         'Three Roles, Layered Permissions',
         size=14, color=TEXT_MUTED, italic=True, font=FONT_EN)

roles = [
    {
        'icon': '🌐', 'color': TEXT_MUTED, 'badge_color': RGBColor(0xE2, 0xE8, 0xF0),
        'zh': '訪客', 'en': 'Public Visitor',
        'features': [
            ('✓', '瀏覽石材庫存', 'Browse inventory'),
            ('✓', '搜尋與篩選', 'Search & filter'),
            ('✓', 'AI 以圖搜圖', 'Image search'),
            ('✓', '用量計算機', 'Quantity calculator'),
            ('🔒', '詢價、設計工具', 'Inquiry, Design Tool'),
        ]
    },
    {
        'icon': '👤', 'color': BLUE, 'badge_color': RGBColor(0xDB, 0xEA, 0xFE),
        'zh': '會員', 'en': 'Member',
        'features': [
            ('✓', '訪客全部功能', 'All visitor features'),
            ('✓', '送出詢價', 'Submit inquiries'),
            ('✓', 'AI 設計合成', 'AI design synthesis'),
            ('✓', '詢問紀錄追蹤', 'Inquiry tracking'),
            ('✓', '個人帳號管理', 'Account management'),
        ]
    },
    {
        'icon': '⚙', 'color': PURPLE, 'badge_color': RGBColor(0xF3, 0xE8, 0xFF),
        'zh': '管理員', 'en': 'Administrator',
        'features': [
            ('✓', '石材 CRUD', 'Stone CRUD'),
            ('✓', '會員管理', 'Member management'),
            ('✓', '詢問狀態維護', 'Inquiry workflow'),
            ('✓', 'Excel 匯出', 'Excel export'),
            ('✓', 'AI Prompt 編輯', 'Prompt editing'),
        ]
    }
]

card_w = 4.0
gap = 0.25
total_w = card_w * 3 + gap * 2
start_x = (13.333 - total_w) / 2
card_y = 2.2

for i, role in enumerate(roles):
    x = start_x + i * (card_w + gap)
    add_rounded(slide, Inches(x), Inches(card_y), Inches(card_w), Inches(4.7),
                WHITE, line=BORDER, line_w=Pt(0.75), radius=0.04)
    # Icon 大圓
    add_oval(slide, Inches(x + (card_w - 0.9) / 2), Inches(card_y + 0.4),
             Inches(0.9), Inches(0.9), role['color'])
    add_text(slide, Inches(x + (card_w - 0.9) / 2), Inches(card_y + 0.48),
             Inches(0.9), Inches(0.75), role['icon'],
             size=28, align='center', color=WHITE, font=FONT_HEAD)
    # 角色名稱
    add_text(slide, Inches(x), Inches(card_y + 1.5), Inches(card_w), Inches(0.4),
             role['zh'], size=20, bold=True, color=TEXT_DARK, align='center', font=FONT_HEAD)
    add_text(slide, Inches(x), Inches(card_y + 1.92), Inches(card_w), Inches(0.3),
             role['en'], size=11, color=role['color'], bold=True, align='center', font=FONT_EN)
    # 分隔線（用細矩形代替）
    add_rect(slide, Inches(x + 0.5), Inches(card_y + 2.35),
             Inches(card_w - 1), Inches(0.015), BORDER)
    # 功能清單
    for j, (mark, zh, en) in enumerate(role['features']):
        y = card_y + 2.55 + j * 0.42
        mark_color = role['color'] if mark == '✓' else TEXT_LIGHT
        add_text(slide, Inches(x + 0.4), Inches(y), Inches(0.3), Inches(0.3),
                 mark, size=11, bold=True, color=mark_color, font=FONT_EN)
        text_color = TEXT_DARK if mark == '✓' else TEXT_LIGHT
        add_text(slide, Inches(x + 0.75), Inches(y), Inches(card_w - 1.2), Inches(0.3),
                 zh, size=11, color=text_color, font=FONT_HEAD)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 8: 技術架構 Tech Stack
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, SW, SH, WHITE)

add_text(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4), 'SECTION 06',
         size=11, bold=True, color=PURPLE, font=FONT_EN)
add_text(slide, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         '技術架構',
         size=32, bold=True, color=TEXT_DARK, font=FONT_HEAD)
add_text(slide, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
         'Tech Stack — Lightweight, Modern, Cloud-Native',
         size=14, color=TEXT_MUTED, italic=True, font=FONT_EN)

tech_groups = [
    ('🐍', '後端 Backend', BLUE, ['Python 3.9', 'Flask 3.1', 'SQLite + WAL', 'Werkzeug Auth']),
    ('🤖', 'AI 服務 AI Services', PURPLE, ['Replicate', 'Flux Fill Pro', 'LLaVA v1.6', 'Pillow + ImageHash']),
    ('☁', '部署 Deployment', BLUE_LIGHT, ['Railway PaaS', 'Gunicorn', 'Volume 持久化', 'GitHub 自動部署']),
    ('🔒', '安全 Security', GREEN, ['PBKDF2-SHA256', 'Session 雙軌', 'CSRF Token', 'TLS 1.3']),
    ('📧', '通訊 Communications', GOLD, ['Flask-Mail', 'Gmail SMTP', 'TLS 加密', 'Email 通知']),
    ('📊', '匯出 Data Export', PINK, ['openpyxl', 'Excel .xlsx', '即時生成', '一鍵下載']),
]

cols = 3; rows = 2
gap = 0.25
card_w = (13.333 - 1.4 - gap * (cols - 1)) / cols
card_h = 2.3
start_x = 0.7
start_y = 2.3

for i, (emoji, title, color, items) in enumerate(tech_groups):
    col = i % cols
    row = i // cols
    x = start_x + col * (card_w + gap)
    y = start_y + row * (card_h + gap)
    add_rounded(slide, Inches(x), Inches(y), Inches(card_w), Inches(card_h),
                WHITE, line=BORDER, line_w=Pt(0.75), radius=0.04)
    # 頂部彩色 badge
    add_rounded(slide, Inches(x + 0.35), Inches(y + 0.3), Inches(0.55), Inches(0.55),
                color, radius=0.5)
    add_text(slide, Inches(x + 0.35), Inches(y + 0.35), Inches(0.55), Inches(0.5),
             emoji, size=18, align='center', color=WHITE, font=FONT_HEAD)
    # 標題
    add_text(slide, Inches(x + 1.05), Inches(y + 0.38), Inches(card_w - 1.4), Inches(0.4),
             title, size=13, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    # 項目清單
    for j, item in enumerate(items):
        item_y = y + 1.1 + j * 0.28
        # 小圓點
        add_oval(slide, Inches(x + 0.4), Inches(item_y + 0.1), Inches(0.08), Inches(0.08), color)
        add_text(slide, Inches(x + 0.6), Inches(item_y), Inches(card_w - 0.9), Inches(0.28),
                 item, size=10.5, color=TEXT_DARK, font=FONT_EN)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 9: 資安合規 Security & Compliance
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, SW, SH, WHITE)

add_text(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4), 'SECTION 07',
         size=11, bold=True, color=PURPLE, font=FONT_EN)
add_text(slide, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         '資安合規與資料保護',
         size=32, bold=True, color=TEXT_DARK, font=FONT_HEAD)
add_text(slide, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
         'Security & Compliance — Aligned with MOEA Guidelines',
         size=14, color=TEXT_MUTED, italic=True, font=FONT_EN)

# 三大區塊
compliance = [
    ('🚫', '排除中資技術', 'No Chinese Tech',
     ['全部 18 項組件已稽核', 'No TikTok / Tencent / Alibaba', 'No .cn 網域引用', '符合行政院限制使用原則']),
    ('🔐', '資料保護', 'Data Protection',
     ['PBKDF2-SHA256 密碼雜湊', 'TLS 1.3 全程加密', 'Session 雙軌隔離', '會員 / 管理員權限分離']),
    ('📋', '法規遵循', 'Regulatory Compliance',
     ['個人資料保護法', '資通安全管理法', '經濟部產業發展署規範', 'GDPR / SOC 2 相容']),
]
card_w = 4.0
gap = 0.25
total_w = card_w * 3 + gap * 2
start_x = (13.333 - total_w) / 2

for i, (emoji, zh, en, items) in enumerate(compliance):
    x = start_x + i * (card_w + gap)
    add_rounded(slide, Inches(x), Inches(2.3), Inches(card_w), Inches(3.8),
                LIGHT_BG, line=BORDER, line_w=Pt(0.75), radius=0.04)
    # Icon
    add_oval(slide, Inches(x + (card_w - 0.8) / 2), Inches(2.55),
             Inches(0.8), Inches(0.8), GREEN)
    add_text(slide, Inches(x + (card_w - 0.8) / 2), Inches(2.63),
             Inches(0.8), Inches(0.65), emoji,
             size=24, align='center', color=WHITE, font=FONT_HEAD)
    # 標題
    add_text(slide, Inches(x), Inches(3.5), Inches(card_w), Inches(0.4),
             zh, size=16, bold=True, color=TEXT_DARK, align='center', font=FONT_HEAD)
    add_text(slide, Inches(x), Inches(3.88), Inches(card_w), Inches(0.3),
             en, size=10, color=GREEN, bold=True, align='center', font=FONT_EN)
    # 項目
    for j, item in enumerate(items):
        item_y = 4.4 + j * 0.36
        add_text(slide, Inches(x + 0.35), Inches(item_y), Inches(0.25), Inches(0.3),
                 '✓', size=11, bold=True, color=GREEN, font=FONT_EN)
        add_text(slide, Inches(x + 0.65), Inches(item_y), Inches(card_w - 0.9), Inches(0.3),
                 item, size=11, color=TEXT_DARK, font=FONT_HEAD)

# 底部三份合規文件提示
add_rounded(slide, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.7),
            NAVY_DEEP, radius=0.04)
add_text(slide, Inches(1.0), Inches(6.6), Inches(11), Inches(0.3),
         '📑 已備齊三份正式合規文件：技術組件聲明書、資安自評表、資料流向圖（共 11 頁）',
         size=11, bold=True, color=WHITE, font=FONT_HEAD)
add_text(slide, Inches(1.0), Inches(6.85), Inches(11), Inches(0.25),
         'Three official compliance documents prepared: Components Declaration, Security Self-Assessment, Data Flow Diagram',
         size=9, color=TEXT_LIGHT, italic=True, font=FONT_EN)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 10: 成本結構 Cost Structure
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, SW, SH, WHITE)

add_text(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4), 'SECTION 08',
         size=11, bold=True, color=PURPLE, font=FONT_EN)
add_text(slide, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         '成本結構：用多少 · 付多少',
         size=32, bold=True, color=TEXT_DARK, font=FONT_HEAD)
add_text(slide, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
         'Usage-Based Pricing — Zero When Idle',
         size=14, color=TEXT_MUTED, italic=True, font=FONT_EN)

# 左側大型 stat 數字
add_rounded(slide, Inches(0.7), Inches(2.3), Inches(4.5), Inches(4.7),
            NAVY_DEEP, radius=0.04)
add_text(slide, Inches(0.9), Inches(2.6), Inches(4.1), Inches(0.4),
         'MONTHLY FIXED COST',
         size=10, bold=True, color=BLUE_LIGHT, font=FONT_EN)
add_text(slide, Inches(0.9), Inches(2.9), Inches(4.1), Inches(0.35),
         '月固定成本',
         size=12, color=TEXT_LIGHT, font=FONT_HEAD)
add_text(slide, Inches(0.9), Inches(3.4), Inches(4.1), Inches(1.5),
         'NT$0',
         size=88, bold=True, color=GREEN, font=FONT_EN)
add_text(slide, Inches(0.9), Inches(5.05), Inches(4.1), Inches(0.35),
         '不需月費 / 沒人用就 0 元',
         size=12, bold=True, color=WHITE, font=FONT_HEAD)
add_text(slide, Inches(0.9), Inches(5.4), Inches(4.1), Inches(0.3),
         'No subscription · No minimum',
         size=10, color=TEXT_LIGHT, italic=True, font=FONT_EN)
# 分隔線
add_rect(slide, Inches(0.9), Inches(5.85), Inches(4.1), Inches(0.015), TEXT_LIGHT)
add_text(slide, Inches(0.9), Inches(6.05), Inches(4.1), Inches(0.3),
         '✓ 完整功能、無限會員、無限石材',
         size=10, color=BLUE_LIGHT, font=FONT_HEAD)
add_text(slide, Inches(0.9), Inches(6.4), Inches(4.1), Inches(0.3),
         '✓ All features, unlimited members & stones',
         size=9, color=TEXT_LIGHT, italic=True, font=FONT_EN)

# 右側：隨用計費表
add_text(slide, Inches(5.7), Inches(2.3), Inches(7), Inches(0.4),
         '💳 隨用計費 / Usage-Based',
         size=14, bold=True, color=PURPLE, font=FONT_HEAD)

usage_items = [
    ('🎨', 'AI 設計合成', 'AI Design Synthesis', '~NT$1.5', '/張 per image'),
    ('👁', 'AI 視覺分析', 'AI Vision Analysis', '~NT$0.5', '/張 per image'),
    ('☁', 'Railway 主機', 'Railway Hosting', 'NT$0', '免費額度內'),
    ('💾', 'Volume 儲存', 'Persistent Storage', '~NT$8', '/GB/月 per GB/month'),
    ('📧', 'Email 通知', 'Email Notifications', 'NT$0', '不額外收費 free'),
]
for i, (emoji, zh, en, price, unit) in enumerate(usage_items):
    y = 2.85 + i * 0.78
    add_rounded(slide, Inches(5.7), Inches(y), Inches(7), Inches(0.7),
                LIGHT_BG, line=BORDER, line_w=Pt(0.5), radius=0.04)
    add_text(slide, Inches(5.9), Inches(y + 0.15), Inches(0.4), Inches(0.4),
             emoji, size=14, color=BLUE, font=FONT_HEAD)
    add_text(slide, Inches(6.4), Inches(y + 0.07), Inches(3.5), Inches(0.35),
             zh, size=12, bold=True, color=TEXT_DARK, font=FONT_HEAD)
    add_text(slide, Inches(6.4), Inches(y + 0.38), Inches(3.5), Inches(0.3),
             en, size=9.5, color=TEXT_MUTED, italic=True, font=FONT_EN)
    add_text(slide, Inches(10), Inches(y + 0.12), Inches(2.6), Inches(0.35),
             price, size=15, bold=True, color=PURPLE, align='right', font=FONT_EN)
    add_text(slide, Inches(10), Inches(y + 0.42), Inches(2.6), Inches(0.25),
             unit, size=9, color=TEXT_MUTED, align='right', italic=True, font=FONT_EN)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 11: 路線圖 Roadmap
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, SW, SH, WHITE)

add_text(slide, Inches(0.7), Inches(0.5), Inches(8), Inches(0.4), 'ROADMAP',
         size=11, bold=True, color=PURPLE, font=FONT_EN)
add_text(slide, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         '產品演進路線圖',
         size=32, bold=True, color=TEXT_DARK, font=FONT_HEAD)
add_text(slide, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
         'Product Evolution Path',
         size=14, color=TEXT_MUTED, italic=True, font=FONT_EN)

roadmap = [
    ('短期', 'SHORT-TERM', '1-3 個月', BLUE, [
        ('SAM 自動辨識區域', 'Auto surface detection'),
        ('合成歷史紀錄', 'Generation history'),
        ('並排比對模式', 'Side-by-side compare'),
        ('行動裝置優化', 'Mobile UX polish'),
    ]),
    ('中期', 'MID-TERM', '3-6 個月', PURPLE, [
        ('Claude / GPT 風格推薦', 'AI style recommendations'),
        ('多角度合成', 'Multi-angle synthesis'),
        ('業主後台分析', 'Vendor analytics'),
        ('白牌授權', 'White-label licensing'),
    ]),
    ('長期', 'LONG-TERM', '6+ 個月', PINK, [
        ('跨材質 SaaS', 'Cross-material SaaS'),
        ('木材 / 皮革 / 磁磚', 'Wood / Leather / Tile'),
        ('東南亞市場拓展', 'SEA market expansion'),
        ('循環經濟生態系', 'Circular economy ecosystem'),
    ]),
]

card_w = 4.0
gap = 0.25
total_w = card_w * 3 + gap * 2
start_x = (13.333 - total_w) / 2

# 連接線
line_y = 3.05
add_rect(slide, Inches(start_x + 1.5), Inches(line_y), Inches(total_w - 3), Inches(0.025), BORDER)

for i, (zh, en, period, color, items) in enumerate(roadmap):
    x = start_x + i * (card_w + gap)
    # 階段圓點
    add_oval(slide, Inches(x + (card_w - 0.5) / 2), Inches(2.8), Inches(0.5), Inches(0.5), color)
    add_text(slide, Inches(x + (card_w - 0.5) / 2), Inches(2.88), Inches(0.5), Inches(0.4),
             str(i + 1), size=14, bold=True, color=WHITE, align='center', font=FONT_EN)
    # 卡片
    add_rounded(slide, Inches(x), Inches(3.5), Inches(card_w), Inches(3.5),
                LIGHT_BG, line=color, line_w=Pt(1), radius=0.04)
    # 中文
    add_text(slide, Inches(x), Inches(3.7), Inches(card_w), Inches(0.4),
             zh, size=18, bold=True, color=TEXT_DARK, align='center', font=FONT_HEAD)
    add_text(slide, Inches(x), Inches(4.12), Inches(card_w), Inches(0.3),
             en, size=10, color=color, bold=True, align='center', font=FONT_EN)
    add_text(slide, Inches(x), Inches(4.42), Inches(card_w), Inches(0.3),
             period, size=10, color=TEXT_MUTED, align='center', italic=True, font=FONT_EN)
    # 項目
    for j, (item_zh, item_en) in enumerate(items):
        item_y = 4.95 + j * 0.5
        add_text(slide, Inches(x + 0.35), Inches(item_y), Inches(0.25), Inches(0.3),
                 '→', size=11, bold=True, color=color, font=FONT_EN)
        add_text(slide, Inches(x + 0.65), Inches(item_y - 0.02), Inches(card_w - 0.9), Inches(0.3),
                 item_zh, size=11, bold=True, color=TEXT_DARK, font=FONT_HEAD)
        add_text(slide, Inches(x + 0.65), Inches(item_y + 0.22), Inches(card_w - 0.9), Inches(0.25),
                 item_en, size=8.5, color=TEXT_MUTED, italic=True, font=FONT_EN)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 12: 結語 Closing
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, SW, SH, NAVY_DEEP)
# 裝飾大鑽石
add_hex(slide, Inches(8), Inches(-3), Inches(10), Inches(10), NAVY_MID)
add_hex(slide, Inches(-3), Inches(5), Inches(7), Inches(7), NAVY_MID)

draw_diamond_logo(slide, 6.65, 2.4, 1.8, BLUE_LIGHT, PURPLE_LIGHT)

# Slogan
add_text(slide, Inches(0.5), Inches(3.7), Inches(12.333), Inches(0.9),
         '讓每一塊石材，找到對的設計',
         size=44, bold=True, color=WHITE, align='center', font=FONT_HEAD)
add_text(slide, Inches(0.5), Inches(4.65), Inches(12.333), Inches(0.5),
         'Let Every Stone Find Its Design',
         size=18, color=BLUE_LIGHT, italic=True, align='center', font=FONT_EN)

# 訪問資訊
add_rounded(slide, Inches(3.5), Inches(5.6), Inches(6.333), Inches(0.7),
            BLUE, radius=0.5)
add_text(slide, Inches(3.5), Inches(5.78), Inches(6.333), Inches(0.5),
         'stone-database-production.up.railway.app',
         size=14, bold=True, color=WHITE, align='center', font='Consolas')

# 標語
add_text(slide, Inches(0.5), Inches(6.6), Inches(12.333), Inches(0.3),
         '— ReStone 循環石材服務平台 · Made for Sustainable Design —',
         size=10, color=TEXT_LIGHT, align='center', italic=True, font=FONT_EN)
add_text(slide, Inches(0.5), Inches(6.95), Inches(12.333), Inches(0.3),
         '2026 · Thank you  ·  謝謝聆聽',
         size=10, color=TEXT_LIGHT, align='center', font=FONT_EN)


# ─── 輸出 ───
output = 'ReStone_簡介簡報.pptx'
prs.save(output)
print(f'✓ 簡報已產出：{output}')
print(f'  共 {len(prs.slides)} 頁')
