"""
ReStone 簡介說明 — 親切版 2 頁簡報
產出：ReStone_簡介說明.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── 親切色系 ───
CREAM       = RGBColor(0xFF, 0xF8, 0xF0)
CREAM_SOFT  = RGBColor(0xFA, 0xF0, 0xE6)
PAPER       = RGBColor(0xFF, 0xFF, 0xFF)
INK         = RGBColor(0x2C, 0x2A, 0x33)
INK_SOFT    = RGBColor(0x6B, 0x68, 0x78)
INK_LIGHT   = RGBColor(0x9B, 0x98, 0xA8)

BLUE        = RGBColor(0x5B, 0x8F, 0xF9)
BLUE_SOFT   = RGBColor(0xD6, 0xE4, 0xFB)
TEAL        = RGBColor(0x5A, 0xC8, 0xB0)
TEAL_SOFT   = RGBColor(0xD1, 0xF2, 0xEA)
PURPLE      = RGBColor(0x9B, 0x8A, 0xFB)
PURPLE_SOFT = RGBColor(0xE3, 0xDE, 0xFD)
PINK        = RGBColor(0xFF, 0x9A, 0xAA)
PINK_SOFT   = RGBColor(0xFF, 0xE0, 0xE6)
YELLOW_SOFT = RGBColor(0xFF, 0xF1, 0xCC)
GREEN_SOFT  = RGBColor(0xD7, 0xF0, 0xD7)
ORANGE      = RGBColor(0xFF, 0x9F, 0x68)
ORANGE_SOFT = RGBColor(0xFF, 0xE0, 0xCC)

DEEP_TEAL   = RGBColor(0x1A, 0x88, 0x70)
DEEP_AMBER  = RGBColor(0x92, 0x58, 0x0A)
DEEP_GREEN  = RGBColor(0x2D, 0x7C, 0x2D)

FONT = 'Microsoft JhengHei'
FONT_NUM = 'Helvetica Neue'

# ─── 16:9 widescreen ───
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


# ─── 工具函式 ───
def add_rect(slide, x, y, w, h, fill, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w: s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_rounded(slide, x, y, w, h, fill, line=None, line_w=None, radius=0.1):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w: s.line.width = line_w
    s.shadow.inherit = False
    try:
        s.adjustments[0] = radius
    except Exception:
        pass
    return s


def add_oval(slide, x, y, w, h, fill, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w: s.line.width = line_w
    s.shadow.inherit = False
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             align='left', valign='top', font=FONT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    valign_map = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = valign_map.get(valign, MSO_ANCHOR.TOP)
    align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
    p = tf.paragraphs[0]
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if font:
        r.font.name = font
    return tb


def draw_diamond_logo(slide, cx_in, cy_in, size_in, fill=BLUE):
    """簡化版鑽石 LOGO（六角形 + 中央菱形）"""
    s = size_in
    x = cx_in - s / 2; y = cy_in - s / 2
    hex_shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(x), Inches(y), Inches(s), Inches(s))
    hex_shape.fill.solid(); hex_shape.fill.fore_color.rgb = fill
    hex_shape.line.fill.background()
    hex_shape.shadow.inherit = False
    hex_shape.rotation = 90
    inner = s * 0.4
    ix = cx_in - inner / 2; iy = cy_in - inner / 2
    d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(ix), Inches(iy), Inches(inner), Inches(inner))
    d.fill.solid(); d.fill.fore_color.rgb = PURPLE
    d.line.fill.background()
    d.shadow.inherit = False


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 1: ReStone 是什麼 + 三個角色
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, SW, SH, CREAM)

# 標題列：品牌 + 一句話
draw_diamond_logo(slide, 1.1, 0.85, 0.7, BLUE)
add_text(slide, Inches(1.55), Inches(0.55), Inches(4), Inches(0.7),
         'ReStone', size=32, bold=True, color=INK, font='Helvetica Neue')
add_text(slide, Inches(1.55), Inches(1.05), Inches(6), Inches(0.4),
         '循環石材服務平台', size=12, color=INK_SOFT)

# 右上角小貼紙
add_rounded(slide, Inches(9.6), Inches(0.65), Inches(3), Inches(0.55), PAPER,
            line=INK_LIGHT, line_w=Pt(1), radius=0.5)
add_text(slide, Inches(9.6), Inches(0.7), Inches(3), Inches(0.5),
         '這個平台給誰用？', size=11, bold=True, color=INK_SOFT,
         align='center', valign='middle')

# 主標
add_text(slide, Inches(0.7), Inches(1.7), Inches(12), Inches(0.7),
         '三種人，原本都有點困擾',
         size=30, bold=True, color=INK)
add_text(slide, Inches(0.7), Inches(2.35), Inches(12), Inches(0.4),
         '讓設計師、廠商、客戶在同一個地方，找到對的石頭。',
         size=13, color=INK_SOFT)

# 三個角色卡片
roles = [
    {
        'bg': PAPER, 'border': ORANGE_SOFT, 'ico_bg': ORANGE_SOFT,
        'emoji': '🏗', 'zh': '石材廠商',
        'pain': '倉庫裡的剩料很美，\n但設計師看不到，賣不掉。',
    },
    {
        'bg': PAPER, 'border': BLUE_SOFT, 'ico_bg': BLUE_SOFT,
        'emoji': '🎨', 'zh': '設計師',
        'pain': '翻 Excel 找料費時，\n客戶又看不懂規格表。',
    },
    {
        'bg': PAPER, 'border': PINK_SOFT, 'ico_bg': PINK_SOFT,
        'emoji': '🏠', 'zh': '客戶 / 屋主',
        'pain': '看著名稱想像不出來，\n不敢決定下單。',
    },
]
card_w = 3.85
gap = 0.3
total_w = card_w * 3 + gap * 2
start_x = (13.333 - total_w) / 2
card_y = 3.2

for i, role in enumerate(roles):
    x = start_x + i * (card_w + gap)
    # 卡片
    add_rounded(slide, Inches(x), Inches(card_y), Inches(card_w), Inches(3.4),
                role['bg'], line=role['border'], line_w=Pt(2), radius=0.06)
    # Icon 圓底
    add_oval(slide, Inches(x + (card_w - 1.1) / 2), Inches(card_y + 0.35),
             Inches(1.1), Inches(1.1), role['ico_bg'])
    add_text(slide, Inches(x), Inches(card_y + 0.5), Inches(card_w), Inches(0.9),
             role['emoji'], size=44, align='center', valign='middle', font=FONT)
    # 角色名
    add_text(slide, Inches(x), Inches(card_y + 1.65), Inches(card_w), Inches(0.5),
             role['zh'], size=18, bold=True, color=INK, align='center')
    # 「他們的煩惱」標籤
    add_text(slide, Inches(x), Inches(card_y + 2.15), Inches(card_w), Inches(0.3),
             '他們的煩惱', size=9, bold=True, color=INK_LIGHT, align='center')
    # 煩惱內容
    add_text(slide, Inches(x + 0.3), Inches(card_y + 2.5), Inches(card_w - 0.6), Inches(0.8),
             role['pain'], size=11, color=INK_SOFT, align='center')

# 底部：頁碼 + 標語
add_text(slide, Inches(0.7), Inches(7.05), Inches(8), Inches(0.3),
         '讓每一塊石材，找到對的設計', size=10, color=INK_LIGHT, italic=True)
add_text(slide, Inches(10), Inches(7.05), Inches(2.8), Inches(0.3),
         '1 / 2', size=10, color=INK_LIGHT, align='right', font='Helvetica Neue')


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 2: ReStone 怎麼幫忙 + 三步驟 + 三大好處
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
add_rect(slide, 0, 0, SW, SH, CREAM)

# 標題列
draw_diamond_logo(slide, 1.1, 0.85, 0.7, BLUE)
add_text(slide, Inches(1.55), Inches(0.55), Inches(4), Inches(0.7),
         'ReStone', size=32, bold=True, color=INK, font='Helvetica Neue')
add_text(slide, Inches(1.55), Inches(1.05), Inches(6), Inches(0.4),
         '循環石材服務平台', size=12, color=INK_SOFT)

add_rounded(slide, Inches(9.6), Inches(0.65), Inches(3), Inches(0.55), PAPER,
            line=INK_LIGHT, line_w=Pt(1), radius=0.5)
add_text(slide, Inches(9.6), Inches(0.7), Inches(3), Inches(0.5),
         '怎麼運作？只要 30 秒', size=11, bold=True, color=INK_SOFT,
         align='center', valign='middle')

# 主標
add_text(slide, Inches(0.7), Inches(1.7), Inches(12), Inches(0.7),
         '三個動作，就有完成圖',
         size=30, bold=True, color=INK)
add_text(slide, Inches(0.7), Inches(2.35), Inches(12), Inches(0.4),
         'AI 幫每個人做最累的事，三方在同一個地方找到對的石頭。',
         size=13, color=INK_SOFT)

# ─── 三步驟流程（橫式） ───
steps = [
    {'num': '1', 'emoji': '📷', 'name': '上傳設計圖', 'desc': '把客戶想看的\n那個空間照片丟上來'},
    {'num': '2', 'emoji': '🖌', 'name': '筆刷塗一塗', 'desc': '用滑鼠圈出\n想換成石材的地方'},
    {'num': '3', 'emoji': '✨', 'name': 'AI 自動合成', 'desc': '挑塊石材按按鈕\n30 秒看到完成圖'},
]
step_w = 3.5
step_gap = 0.5
total_step_w = step_w * 3 + step_gap * 2
step_start_x = (13.333 - total_step_w) / 2
step_y = 3.0

for i, st in enumerate(steps):
    x = step_start_x + i * (step_w + step_gap)
    # 卡片
    add_rounded(slide, Inches(x), Inches(step_y), Inches(step_w), Inches(2.0),
                PAPER, line=BLUE_SOFT, line_w=Pt(1.5), radius=0.08)
    # 編號圓
    add_oval(slide, Inches(x + 0.3), Inches(step_y + 0.25),
             Inches(0.55), Inches(0.55), BLUE)
    add_text(slide, Inches(x + 0.3), Inches(step_y + 0.28), Inches(0.55), Inches(0.5),
             st['num'], size=14, bold=True, color=PAPER, align='center', font='Helvetica Neue')
    # Emoji
    add_text(slide, Inches(x + 1.0), Inches(step_y + 0.25), Inches(step_w - 1.3), Inches(0.6),
             st['emoji'], size=26, align='left', font=FONT)
    # 名稱
    add_text(slide, Inches(x + 0.3), Inches(step_y + 0.95), Inches(step_w - 0.6), Inches(0.4),
             st['name'], size=15, bold=True, color=INK)
    # 描述
    add_text(slide, Inches(x + 0.3), Inches(step_y + 1.35), Inches(step_w - 0.6), Inches(0.6),
             st['desc'], size=10, color=INK_SOFT)

    # 箭頭（除了最後一張）
    if i < 2:
        arrow_x = x + step_w + (step_gap - 0.3) / 2
        add_text(slide, Inches(arrow_x), Inches(step_y + 0.7),
                 Inches(0.3), Inches(0.6),
                 '→', size=28, color=INK_LIGHT, align='center', valign='middle', font='Helvetica Neue')

# ─── 三大好處 ───
benefits = [
    {'bg': TEAL_SOFT, 'num': '60×', 'num_color': DEEP_TEAL, 'emoji': '⏱',
     'name': '提案更快', 'desc': '原本 30 分鐘 · 現在 30 秒'},
    {'bg': YELLOW_SOFT, 'num': 'NT$0', 'num_color': DEEP_AMBER, 'emoji': '💰',
     'name': '月費全免', 'desc': '沒人用就 0 元 · 用多少付多少'},
    {'bg': GREEN_SOFT, 'num': '3 成', 'num_color': DEEP_GREEN, 'emoji': '🌱',
     'name': '剩料再生', 'desc': '本來會被淘汰的 · 變成新選擇'},
]
ben_w = 3.85
ben_gap = 0.3
total_ben_w = ben_w * 3 + ben_gap * 2
ben_start_x = (13.333 - total_ben_w) / 2
ben_y = 5.25

# 區段標籤
add_text(slide, Inches(0.7), Inches(5.0), Inches(8), Inches(0.3),
         '這樣做有什麼好？', size=11, bold=True, color=INK_SOFT)

for i, b in enumerate(benefits):
    x = ben_start_x + i * (ben_w + ben_gap)
    add_rounded(slide, Inches(x), Inches(ben_y), Inches(ben_w), Inches(1.6),
                b['bg'], radius=0.08)
    # Emoji
    add_text(slide, Inches(x + 0.4), Inches(ben_y + 0.3), Inches(0.7), Inches(0.6),
             b['emoji'], size=24, font=FONT)
    # 大數字
    add_text(slide, Inches(x + 1.2), Inches(ben_y + 0.2), Inches(ben_w - 1.4), Inches(0.7),
             b['num'], size=26, bold=True, color=b['num_color'], font='Helvetica Neue')
    # 名稱
    add_text(slide, Inches(x + 0.4), Inches(ben_y + 0.95), Inches(ben_w - 0.6), Inches(0.4),
             b['name'], size=13, bold=True, color=INK)
    # 描述
    add_text(slide, Inches(x + 0.4), Inches(ben_y + 1.25), Inches(ben_w - 0.6), Inches(0.3),
             b['desc'], size=9, color=INK_SOFT)

# 底部：標語 + URL + 頁碼
add_text(slide, Inches(0.7), Inches(7.05), Inches(8), Inches(0.3),
         '讓每一塊石材，找到對的設計', size=10, color=INK_LIGHT, italic=True)
add_text(slide, Inches(10), Inches(7.05), Inches(2.8), Inches(0.3),
         '2 / 2', size=10, color=INK_LIGHT, align='right', font='Helvetica Neue')


# ─── 輸出 ───
output = 'ReStone_簡介說明.pptx'
prs.save(output)
print(f'✓ 簡報已產出：{output}（共 {len(prs.slides)} 頁）')
