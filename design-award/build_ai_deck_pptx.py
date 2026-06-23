# -*- coding: utf-8 -*-
"""產生「ReStone AI 設計功能與規格」PowerPoint 簡報（深色科技風，16:9）。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FONT = 'Noto Sans TC'
BG     = RGBColor(0x0C, 0x12, 0x1F)
SURF   = RGBColor(0x18, 0x20, 0x33)
TEXT   = RGBColor(0xF2, 0xF6, 0xFF)
MUTED  = RGBColor(0xA8, 0xB6, 0xCC)
CYAN   = RGBColor(0x22, 0xD3, 0xEE)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
GREEN  = RGBColor(0x34, 0xD3, 0x99)
BORDER = RGBColor(0x2A, 0x35, 0x4D)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)
    return s


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    b = s.shapes.add_textbox(x, y, w, h); tf = b.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0); tf.margin_top = tf.margin_bottom = Pt(0)
    return tf


def setp(p, text, size, color, bold=False, align=PP_ALIGN.LEFT, space=6):
    p.text = text; p.alignment = align; p.space_after = Pt(space)
    r = p.runs[0]; r.font.name = r.font.name = FONT
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r._r.rPr.set('lang', 'zh-TW')
    return p


def kicker(s, text):
    setp(tb(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.4)).paragraphs[0],
         text, 13, CYAN, bold=True)


def title(s, text, y=1.15, size=34, color=TEXT):
    tf = tb(s, Inches(0.9), Inches(y), Inches(11.5), Inches(1.6))
    setp(tf.paragraphs[0], text, size, color, bold=True)


def card(s, x, y, w, h, head, headcolor, lines, headsize=15, body=11):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = SURF
    box.line.color.rgb = BORDER; box.line.width = Pt(1); box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.22); tf.margin_top = tf.margin_bottom = Inches(0.18)
    setp(tf.paragraphs[0], head, headsize, headcolor, bold=True, space=8)
    for ln in lines:
        p = tf.add_paragraph(); setp(p, ln, body, MUTED, space=5)
    return box


def bullets(tf, items, size=14, color=MUTED):
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        setp(p, '・' + it, size, color, space=8)


# ── 1 封面 ──
s = slide()
kicker(s, 'RESTONE 循環石材服務平台')
tf = tb(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(2.2))
setp(tf.paragraphs[0], 'AI 設計功能', 54, TEXT, bold=True, space=2)
setp(tf.add_paragraph(), '與技術規格', 54, CYAN, bold=True)
setp(tb(s, Inches(0.9), Inches(4.6), Inches(10.5), Inches(1.5)).paragraphs[0],
     '一個平台，整合多項 AI 設計工具，讓設計師把循環石材即時帶進設計、加速完成高值化產品。', 18, MUTED)

# ── 2 總覽 ──
s = slide(); kicker(s, 'OVERVIEW · 功能總覽'); title(s, '四項 AI 設計工具 + 輔助引擎')
cw, gap, x0, y0, ch = Inches(2.85), Inches(0.22), Inches(0.9), Inches(2.5), Inches(2.4)
data = [('核心 · 材質模擬', CYAN, ['把循環石材真實紋理', '套用到設計圖指定區域']),
        ('核心 · 牆面模擬', CYAN, ['點選四角，正確透視', '與比例即時貼上石材']),
        ('核心 · 草圖渲染', CYAN, ['手繪草圖 + 石材', '→ 寫實產品渲染圖']),
        ('輔助 · 辨識/搜圖', GREEN, ['AI 看石材寫描述', '以圖找相似石材'])]
for i, (h, c, ls) in enumerate(data):
    card(s, x0 + i * (cw + gap), y0, cw, ch, h, c, ls)
setp(tb(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1)).paragraphs[0],
     '加值串接：減碳效益計算、設計提案 PDF、我的設計歷史、後台展覽圖庫。', 15, MUTED)

# ── 3-5 三大核心 ──
core = [
    ('FEATURE 01 · 核心', '材質模擬　Material Simulation',
     '上傳設計圖或空間照片，圈選要替換的區域，AI 把所選循環石材的真實材質貼上去，即時預覽。',
     ['筆刷塗選 + 拖放上傳（支援 HEIC）', '主力：Gemini 參考圖編輯，以實際石材照片為依據',
      '後備：Flux Fill Pro 文字 inpainting', '結果存入「我的設計」、可出提案 PDF'],
     ['模型：gemini-2.5-flash-image', '後備：flux-fill-pro', '成本：約 NT$1.2 / 張', '權限：僅限登入會員']),
    ('FEATURE 02 · 核心', '牆面模擬　Wall Simulator',
     '點選牆面四角、輸入實際尺寸，以幾何運算精準貼合——透視與比例皆數學精準，並同步顯示用量、庫存與減碳。',
     ['四角點選（支援觸控、可拖曳微調）', '透視投影 + 依實際尺寸算拼貼片數',
      '棋盤鏡像對花、可選拼縫線、光影轉移', '同步：拼貼計畫 / 庫存檢查 / 減碳效益'],
     ['引擎：本地幾何運算（Pillow + NumPy）', '成本：免費・即時出圖', '選配：Gemini AI 擬真強化', '適用：平整牆面 / 桌面 / 地板']),
    ('FEATURE 03 · 核心', '草圖渲染　Sketch Render',
     '上傳手繪草圖或概念線稿，選一塊循環石材，補一句描述，AI 生成寫實產品渲染圖——概念秒變提案。',
     ['草圖 + 石材 + 文字描述', '免遮罩：整張草圖即產品造型',
      '輸出工作室等級渲染（可當型錄主圖）', '自動存入「我的設計」、可出提案 PDF'],
     ['模型：gemini-2.5-flash-image', '成本：約 NT$1.2 / 張', '輸入：JPG / PNG / HEIC / WebP', '權限：僅限登入會員']),
]
for kik, ttl, lead, feat, spec in core:
    s = slide(); kicker(s, kik); title(s, ttl, size=30)
    setp(tb(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.1)).paragraphs[0], lead, 15, MUTED)
    c1 = card(s, Inches(0.9), Inches(3.25), Inches(5.7), Inches(3.4), '功能', CYAN, [])
    bullets(c1.text_frame, feat)
    c2 = card(s, Inches(6.85), Inches(3.25), Inches(5.6), Inches(3.4), '規格', GREEN, [])
    bullets(c2.text_frame, spec)

# ── 6 輔助引擎 ──
s = slide(); kicker(s, 'SUPPORTING · 輔助引擎'); title(s, '讓設計更聰明的底層 AI')
sup = [('視覺辨識', ['AI「看」石材照片自動產生', '精準材質描述並快取重用', '模型：gemini-2.5-flash', '後備：LLaVA v1.6']),
       ('以圖搜圖', ['上傳參考圖，pHash 感知雜湊', '+ HSV 色彩比對找相似石材', '本地運算・免費', 'ImageHash + 色彩分析']),
       ('三層 Prompt 工程', ['管理員自訂 ＞ 視覺 AI 描述', '＞ 模板後備', '自動快取至資料庫', '設計師可即時微調'])]
cw3, x0 = Inches(3.85), Inches(0.9)
for i, (h, ls) in enumerate(sup):
    card(s, x0 + i * (cw3 + Inches(0.22)), Inches(2.5), cw3, Inches(3.6), h, CYAN, ls)

# ── 7 加值 ──
s = slide(); kicker(s, 'VALUE-ADD · 加值串接'); title(s, '從「看效果」到「可成案」')
val = [('減碳效益', ['依體積 × ICE Database', 'v3.0 碳係數自動計算', '等效樹／車換算']),
       ('提案 PDF', ['一鍵產出 A4 提案', '效果圖+規格+減碳+供應商', '內嵌中文字型']),
       ('我的設計', ['每位會員自動保留', '最近 30 張合成作品', '回顧設計發展']),
       ('展覽圖庫', ['後台 Gemini 批次生成', '石材 × 情境展示圖', '供展覽使用'])]
cw4 = Inches(2.85)
for i, (h, ls) in enumerate(val):
    card(s, Inches(0.9) + i * (cw4 + Inches(0.22)), Inches(2.5), cw4, Inches(3.4), h, GREEN, ls)

# ── 8 規格表 ──
s = slide(); kicker(s, 'ARCHITECTURE · 技術規格'); title(s, '多引擎架構與規格')
rows = [('項目', '規格 / 模型', '成本'),
        ('材質模擬（主力）', 'Gemini 2.5 Flash Image 參考圖編輯', '~NT$1.2/張'),
        ('材質模擬（後備）', 'Flux Fill Pro（Replicate）', '~NT$1.6/張'),
        ('草圖渲染', 'Gemini 2.5 Flash Image', '~NT$1.2/張'),
        ('牆面模擬', '本地透視幾何運算（Pillow+NumPy）', '免費・即時'),
        ('視覺辨識', 'Gemini 2.5 Flash（後備 LLaVA）', '極低/免費'),
        ('以圖搜圖', 'pHash 感知雜湊 + HSV 色彩比對', '免費・本地'),
        ('部署/儲存', 'Flask+Gunicorn · Railway · SQLite+Volume', '—')]
tbl = s.shapes.add_table(len(rows), 3, Inches(0.9), Inches(2.3),
                         Inches(11.5), Inches(4.6)).table
tbl.columns[0].width = Inches(3.0); tbl.columns[1].width = Inches(6.3); tbl.columns[2].width = Inches(2.2)
for ri, row in enumerate(rows):
    for ci, val_ in enumerate(row):
        cell = tbl.cell(ri, ci); cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x14, 0x1C, 0x2E) if ri else RGBColor(0x10, 0x2A, 0x33)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.12)
        cell.margin_top = cell.margin_bottom = Inches(0.05)
        p = cell.text_frame.paragraphs[0]
        setp(p, val_, 12 if ri else 11, (CYAN if ri == 0 else (TEXT if ci == 0 else MUTED)),
             bold=(ri == 0 or ci == 0), space=0)

# ── 9 策略 ──
s = slide(); kicker(s, 'STRATEGY · 成本與合規'); title(s, '務實的資源控管')
strat = [('4 / 7', GREEN, ['逾半數引擎為本地或免費運算', '（牆面模擬、以圖搜圖、色彩偵測）']),
         ('會員制', CYAN, ['AI 工具僅限登入會員', '保護運算資源、避免成本暴衝']),
         ('無中資', PURPLE, ['技術堆疊皆非中資來源', 'Google / BFL / 開源，合規'])]
cw3 = Inches(3.85)
for i, (n, c, ls) in enumerate(strat):
    box = card(s, Inches(0.9) + i * (cw3 + Inches(0.22)), Inches(2.5), cw3, Inches(3.0), n, c, ls, headsize=30, body=13)
setp(tb(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1)).paragraphs[0],
     '主力以 Gemini，失敗自動退回 Flux；牆面模擬走本地運算。穩定、低成本、可隨 AI 技術演進同步升級。', 14, MUTED)

# ── 10 結語 ──
s = slide(); kicker(s, 'SUMMARY')
tf = tb(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.2))
setp(tf.paragraphs[0], 'AI × 循環石材', 46, TEXT, bold=True, space=2)
setp(tf.add_paragraph(), '加速高值化設計', 46, CYAN, bold=True)
setp(tb(s, Inches(0.9), Inches(4.8), Inches(11), Inches(1.4)).paragraphs[0],
     '把「找到 → 設計 → 成案」一氣呵成——讓每一塊循環石材，都能再次被好好使用。', 17, MUTED)
setp(tb(s, Inches(0.9), Inches(6.6), Inches(11), Inches(0.5)).paragraphs[0],
     'ReStone 循環石材服務平台　·　支持循環經濟，讓石材物盡其用', 12, MUTED)

out = 'design-award/ReStone_AI設計功能與規格.pptx'
prs.save(out)
print('saved', out, '| slides:', len(prs.slides._sldIdLst))
